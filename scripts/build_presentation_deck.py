"""
Script to generate an executive presentation deck (PPTX) for SIGMA-K Design & Prototype.
Output: docs/PRESENTASI_PERANCANGAN_SIGMA-K_v1.0.pptx
"""

import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PPTX_OUT = os.path.abspath('docs/PRESENTASI_PERANCANGAN_SIGMA-K_v1.0.pptx')
ASSETS_DIR = os.path.abspath('docs/assets')

# Colors
C_NAVY = RGBColor(11, 42, 74)       # #0B2A4A
C_BLUE = RGBColor(30, 64, 175)     # #1E40AF
C_GOLD = RGBColor(212, 175, 55)    # #D4AF37
C_SLATE = RGBColor(30, 41, 59)     # #1E293B
C_MUTED = RGBColor(100, 116, 139)  # #64748B
C_WHITE = RGBColor(255, 255, 255)
C_BG_LIGHT = RGBColor(248, 250, 252)

def add_header(slide, title_text, category="E-SKLD / SIGMA-K — KEMENPANRB"):
    # Header Banner
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_NAVY
    top_bar.line.color.rgb = C_NAVY

    # Gold Accent Line
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = C_GOLD
    gold_line.line.color.rgb = C_GOLD

    # Category Text
    tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.7), Inches(0.3))
    p_cat = tx_cat.text_frame.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_GOLD

    # Title Text
    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.6))
    p_title = tx_title.text_frame.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

def add_footer(slide):
    tx_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3))
    p_foot = tx_foot.text_frame.paragraphs[0]
    p_foot.text = "Kementerian Pendayagunaan Aparatur Negara dan Reformasi Birokrasi (KemenPANRB) — Dokumen Perancangan Sistem v1.0"
    p_foot.font.size = Pt(8.5)
    p_foot.font.color.rgb = C_MUTED

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================
    # SLIDE 1: COVER
    # =========================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_NAVY
    bg1.line.color.rgb = C_NAVY

    # Gold Accent Box
    gb = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.15), Inches(4.8))
    gb.fill.solid()
    gb.fill.fore_color.rgb = C_GOLD
    gb.line.color.rgb = C_GOLD

    tx1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(4.8))
    tf1 = tx1.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "KEMENTERIAN PENDAYAGUNAAN APARATUR NEGARA DAN REFORMASI BIROKRASI"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = C_GOLD

    p1 = tf1.add_paragraph()
    p1.text = "DOKUMEN PERANCANGAN SISTEM &\nPROTOTYPE APLIKASI E-SKLD (SIGMA-K)"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.space_before = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "Sistem Pengelolaan Data Kementerian/Lembaga/Pemerintah Daerah dan Struktur Kelembagaan"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(191, 219, 254)
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Versi 1.0  |  Status: 100% Lolos Pengujian Otomatis  |  27 Agustus 2026"
    p3.font.size = Pt(11)
    p3.font.color.rgb = C_MUTED
    p3.space_before = Pt(30)

    # =========================================================
    # SLIDE 2: EXECUTIVE SUMMARY & OBJECTIVES
    # =========================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. Gambaran Umum & Tujuan Pengembangan SIGMA-K")
    add_footer(s2)

    tx2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Latar Belakang & Urgensi:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_NAVY

    bullets = [
        "Penataan struktur organisasi K/L/D pasca pembentukan Kabinet Merah Putih.",
        "Kebutuhan single source of truth untuk master data instansi, unit, dan jabatan struktural.",
        "Standardisasi alur pengusulan (submission) penataan kelembagaan secara transparan.",
        "Pemisahan wewenang tegas (Separation of Duties): Penapisan Admin (Gate 1) dan Telaah Substantif Verifikator (Gate 2).",
        "Wewenang pengesahan Surat Keputusan (SK) mutlak di tangan Verifikator."
    ]
    for b in bullets:
        bp = tf2.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(11)
        bp.font.color.rgb = C_SLATE
        bp.space_before = Pt(6)

    # Right Card: Validated Metrics
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = C_BG_LIGHT
    card.line.color.rgb = RGBColor(203, 213, 225)

    tx_card = s2.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.1), Inches(4.9))
    tfc = tx_card.text_frame
    tfc.word_wrap = True
    pc = tfc.paragraphs[0]
    pc.text = "Capaian Validasi Sistem (100% PASS):"
    pc.font.size = Pt(13)
    pc.font.bold = True
    pc.font.color.rgb = C_BLUE

    metrics = [
        ("105 Frontend Tests PASS", "Integrasi HTTP, Auth, Security, Master Data & Reports."),
        ("198 Backend PHPUnit Tests", "713 assertions, 0 errors, 0 failures."),
        ("0 TypeScript & ESLint Errors", "Kompilasi strict type-checking 100% bersih."),
        ("16 Next.js Routes Compiled", "Semua rute statis & dinamis terkompilasi sukses."),
        ("21 Relational Database Tables", "Integritas skema MySQL eskld_db terjaga mutlak.")
    ]
    for m_title, m_desc in metrics:
        p_m = tfc.add_paragraph()
        p_m.text = f"✔ {m_title}: {m_desc}"
        p_m.font.size = Pt(10)
        p_m.font.color.rgb = C_SLATE
        p_m.space_before = Pt(8)

    # =========================================================
    # SLIDE 3: SYSTEM ARCHITECTURE
    # =========================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. Arsitektur Sistem 4-Tier E-SKLD / SIGMA-K")
    add_footer(s3)

    img_arch = os.path.join(ASSETS_DIR, 'system_architecture.png')
    if os.path.exists(img_arch):
        s3.shapes.add_picture(img_arch, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 4: ROLES & ZERO TRUST
    # =========================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. Model Otorisasi, Peran Pengguna & Zero-Trust")
    add_footer(s4)

    img_role = os.path.join(ASSETS_DIR, 'role_access_matrix.png')
    if os.path.exists(img_role):
        s4.shapes.add_picture(img_role, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 5: SITEMAP & NAVIGATION
    # =========================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. Struktur Navigasi & Sitemap Aplikasi (16 Rute)")
    add_footer(s5)

    img_site = os.path.join(ASSETS_DIR, 'sitemap_diagram.png')
    if os.path.exists(img_site):
        s5.shapes.add_picture(img_site, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 6: DATABASE ERD
    # =========================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. Struktur Data & Entity Relationship Diagram (21 Tabel)")
    add_footer(s6)

    img_erd = os.path.join(ASSETS_DIR, 'erd_diagram.png')
    if os.path.exists(img_erd):
        s6.shapes.add_picture(img_erd, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 7: MASTER DATA & HIERARCHY
    # =========================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. Model Hierarki Organisasi & Bagan React Flow")
    add_footer(s7)

    img_tree = os.path.join(ASSETS_DIR, 'org_hierarchy_tree.png')
    if os.path.exists(img_tree):
        s7.shapes.add_picture(img_tree, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 8: SUBMISSION FINITE STATE MACHINE
    # =========================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "7. Siklus Hidup Usulan & Finite State Machine (FSM)")
    add_footer(s8)

    img_fsm = os.path.join(ASSETS_DIR, 'submission_lifecycle_fsm.png')
    if os.path.exists(img_fsm):
        s8.shapes.add_picture(img_fsm, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 9: GATE 1 & GATE 2 FLOWCHARTS
    # =========================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "8. Alur Kerja Penapisan Gate 1 & Telaah Substantif Gate 2")
    add_footer(s9)

    img_g1 = os.path.join(ASSETS_DIR, 'gate1_admin_flowchart.png')
    img_g2 = os.path.join(ASSETS_DIR, 'gate2_verifier_flowchart.png')
    if os.path.exists(img_g1) and os.path.exists(img_g2):
        s9.shapes.add_picture(img_g1, Inches(0.6), Inches(1.5), width=Inches(5.8))
        s9.shapes.add_picture(img_g2, Inches(6.8), Inches(1.5), width=Inches(5.8))

    # =========================================================
    # SLIDE 10: VERSIONING & DIFF VIEWER
    # =========================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "9. Snapshot Versioning & Pelacakan Perubahan (Diff Flow)")
    add_footer(s10)

    img_diff = os.path.join(ASSETS_DIR, 'versioning_diff_flow.png')
    if os.path.exists(img_diff):
        s10.shapes.add_picture(img_diff, Inches(0.8), Inches(1.4), width=Inches(11.733))

    # =========================================================
    # SLIDE 11: PROTOTYPE UI: DASHBOARD & KATALOG
    # =========================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "10. Prototype Antarmuka: Dashboard Eksekutif & Master Instansi")
    add_footer(s11)

    img_dash = os.path.join(ASSETS_DIR, 'prototype_dashboard.png')
    img_inst = os.path.join(ASSETS_DIR, 'prototype_institutions.png')
    if os.path.exists(img_dash) and os.path.exists(img_inst):
        s11.shapes.add_picture(img_dash, Inches(0.6), Inches(1.5), width=Inches(5.8))
        s11.shapes.add_picture(img_inst, Inches(6.8), Inches(1.5), width=Inches(5.8))

    # =========================================================
    # SLIDE 12: PROTOTYPE UI: ORGANISASI & USULAN
    # =========================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "11. Prototype Antarmuka: Bagan React Flow & Diff Usulan")
    add_footer(s12)

    img_org = os.path.join(ASSETS_DIR, 'prototype_org_structure.png')
    img_sub = os.path.join(ASSETS_DIR, 'prototype_submission_detail.png')
    if os.path.exists(img_org) and os.path.exists(img_sub):
        s12.shapes.add_picture(img_org, Inches(0.6), Inches(1.5), width=Inches(5.8))
        s12.shapes.add_picture(img_sub, Inches(6.8), Inches(1.5), width=Inches(5.8))

    # =========================================================
    # SLIDE 13: PROTOTYPE UI: VERIFIER & ANALYTICS
    # =========================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "12. Prototype Antarmuka: Ruang Kerja Verifikator & Analitik")
    add_footer(s13)

    img_ver = os.path.join(ASSETS_DIR, 'prototype_verifier_workspace.png')
    img_ana = os.path.join(ASSETS_DIR, 'prototype_analytics_reporting.png')
    if os.path.exists(img_ver) and os.path.exists(img_ana):
        s13.shapes.add_picture(img_ver, Inches(0.6), Inches(1.5), width=Inches(5.8))
        s13.shapes.add_picture(img_ana, Inches(6.8), Inches(1.5), width=Inches(5.8))

    # =========================================================
    # SLIDE 14: KESIMPULAN & PENUTUP
    # =========================================================
    s14 = prs.slides.add_slide(blank_layout)
    bg14 = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg14.fill.solid()
    bg14.fill.fore_color.rgb = C_NAVY
    bg14.line.color.rgb = C_NAVY

    tx14 = s14.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(4.5))
    tf14 = tx14.text_frame
    tf14.word_wrap = True

    p = tf14.paragraphs[0]
    p.text = "KESIMPULAN & KESIAPAN IMPLEMENTASI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    recap = [
        "Arsitektur teruji dan siap produksi (CodeIgniter 4 + MySQL + Next.js 14).",
        "Pemisahan wewenang mutlak (Gate 1 Screening Admin & Gate 2 Final Approval Verifier).",
        "Keamanan Zero-Trust & proteksi BOLA/IDOR terintegrasi di seluruh controller.",
        "Otomasi promosi data pasca pengesahan SK resmi menjamin konsistensi master data nasional.",
        "Seluruh dokumen teknis (Word, PDF, Presentasi PPTX) telah siap diserahkan kepada mentor."
    ]
    for r in recap:
        pr = tf14.add_paragraph()
        pr.text = "✔  " + r
        pr.font.size = Pt(13)
        pr.font.color.rgb = C_WHITE
        pr.space_before = Pt(12)

    print(f"Saving Presentation to {PPTX_OUT}...")
    prs.save(PPTX_OUT)
    print("Presentation successfully created!")

if __name__ == '__main__':
    build_presentation()
